"""
parsers.py
==========
Módulo de ingestão multifonte para o Dashboard Financeiro.

Sistema de detecção automática com 3 estratégias em cascata:

  ESTRATÉGIA 1 — Semicolon-delimited single cell
    Detecta quando o arquivo tem uma única coluna cujos valores contêm
    múltiplos campos separados por ';'. Expande automaticamente.
    Ex: CAP_ILHAS_DO_LAGO_CONCILIADO.xlsx

  ESTRATÉGIA 2 — Named header detection
    Varre todas as linhas procurando uma que contenha palavras-chave
    de cabeçalho (data, valor, descrição...). Extrai múltiplos blocos.
    Ex: JUNHO_2026.xlsx

  ESTRATÉGIA 3 — Headerless positional extraction
    Para arquivos sem cabeçalho nenhum. Detecta a coluna de data por
    padrão de conteúdo (regex), valor por conteúdo numérico/monetário,
    e categoria pelo campo de texto mais longo.
    Ex: CAP_PLAN_IMPORT_CONS_122023_ILA.csv, CAR_PLAN_IMPORT_CONS_052026_HP.csv
"""

import io
import logging
import re
from datetime import date, datetime
from typing import Any, BinaryIO

import pandas as pd

from contaview.logic.leitor_xml_legado import detectar_xml_legado, ler_xml_legado
from contaview.logic.mapeamento_colunas import mapear_colunas, derivar_valor_tipo_de_debito_credito

logger = logging.getLogger(__name__)

COLUNAS_PADRAO: list[str] = [
    "data", "conta_contabil", "valor", "tipo", "historico", "filial"
]
SENTINEL: str = "NÃO ENCONTRADO"

_RE_VALOR = re.compile(r"R?\$?\s*([\d]{1,3}(?:[.,]\d{3})*(?:[.,]\d{1,2})?)")
_RE_DATA  = re.compile(r"\b(\d{2}[/\-]\d{2}[/\-]\d{2,4})\b")
_RE_DATA_ISO = re.compile(r"\b(\d{4}-\d{2}-\d{2})\b")

# Mapeamento de meses em português para números (resolução de datas ambíguas)
MESES_PT: dict[str, int] = {
    'janeiro': 1, 'fevereiro': 2, 'março': 3, 'abril': 4,
    'maio': 5, 'junho': 6, 'julho': 7, 'agosto': 8,
    'setembro': 9, 'outubro': 10, 'novembro': 11, 'dezembro': 12,
    'jan': 1, 'fev': 2, 'mar': 3, 'abr': 4,
    'mai': 5, 'jun': 6, 'jul': 7, 'ago': 8,
    'set': 9, 'out': 10, 'nov': 11, 'dez': 12,
}

# Regex para extrair período do nome do arquivo
_RE_PERIODO_NOME_MM_AAAA = re.compile(r'(\d{2})[-_.](\d{4})')
_RE_PERIODO_NOME_AAAA_MM = re.compile(r'(\d{4})[-_.](\d{2})')
_RE_PERIODO_NOME_MES_ANO = re.compile(
    r'(janeiro|fevereiro|março|abril|maio|junho|julho|agosto|'
    r'setembro|outubro|novembro|dezembro|'
    r'jan|fev|mar|abr|mai|jun|jul|ago|set|out|nov|dez)'
    r'[-_.\s]*?(\d{4})',
    re.IGNORECASE,
)

# Keywords para detecção de linha de cabeçalho (Estratégia 2)
_HEADER_KEYWORDS: dict[str, list[str]] = {
    "data":          ["data", "vencimento", "date", "dt", "competencia", "competência"],
    "conta_contabil": ["conta contábil", "conta", "categoria", "descrição", "descricao",
                      "description", "nome", "historico", "histórico"],
    "valor":         ["valor", "value", "amount", "total", "r$", "preço", "preco",
                      "custo", "saída", "saida", "entrada"],
}

# Mapeamento de variações → campo padrão (Estratégia 2)
_COL_ALIASES: dict[str, list[str]] = {
    "data":          ["data", "vencimento", "date", "dt", "dia", "competencia", "competência"],
    "conta_contabil": ["conta contábil", "conta_contabil", "conta", "categoria",
                      "descrição", "descricao", "description", "nome", "classificação", "classificacao"],
    "valor":         ["valor (r$)", "valor", "saída (r$)", "saida (r$)", "entrada (r$)",
                      "value", "amount", "total (r$)", "total", "vlr", "vl"],
    "tipo":          ["tipo", "tp", "c/d", "type", "nat", "natureza"],
    "historico":     ["histórico", "historico", "hist", "descricao", "descrição",
                      "complemento", "observação", "observacao", "obs"],
    "filial":        ["filial", "fil", "unidade", "cc", "centro de custo"],
}


# ---------------------------------------------------------------------------
# Helpers compartilhados
# ---------------------------------------------------------------------------
def _df_vazio() -> pd.DataFrame:
    return pd.DataFrame(columns=COLUNAS_PADRAO)


def _normalizar_valor(raw: Any) -> float:
    if isinstance(raw, (int, float)):
        return round(float(raw), 2)
    texto = str(raw).strip().replace("R$", "").replace(" ", "")
    if not texto or texto in ("nan", SENTINEL):
        return 0.0
    if "," in texto and "." in texto:
        if texto.rfind(",") > texto.rfind("."):
            texto = texto.replace(".", "").replace(",", ".")
        else:
            texto = texto.replace(",", "")
    elif "," in texto:
        texto = texto.replace(",", ".")
    try:
        return round(float(texto), 2)
    except ValueError:
        return 0.0


def _normalizar_data(raw: Any) -> str:
    if raw is None:
        return SENTINEL
    try:
        if pd.isna(raw):
            return SENTINEL
    except (TypeError, ValueError):
        pass
    if isinstance(raw, (date, datetime)):
        return raw.strftime("%Y-%m-%d") if isinstance(raw, datetime) else raw.isoformat()
    texto = str(raw).strip()
    for fmt in (
        "%Y-%m-%d %H:%M:%S", "%Y-%m-%dT%H:%M:%S",
        "%d/%m/%Y", "%d-%m-%Y", "%Y-%m-%d",
        "%d/%m/%y", "%m/%d/%Y",
    ):
        try:
            return datetime.strptime(texto, fmt).strftime("%Y-%m-%d")
        except ValueError:
            continue
    return SENTINEL


def _resolver_coluna(headers: list[str], campo: str) -> int | None:
    """Resolve índice de coluna por match exato e depois parcial."""
    aliases = _COL_ALIASES.get(campo, [])
    # Exato primeiro
    for i, h in enumerate(headers):
        if str(h).lower().strip() in aliases:
            return i
    # Parcial como fallback
    for i, h in enumerate(headers):
        h_l = str(h).lower().strip()
        if any(a in h_l for a in aliases):
            return i
    return None


def _e_linha_cabecalho(row: pd.Series) -> bool:
    """Retorna True se a linha parece ser um cabeçalho (≥2 categorias de keywords)."""
    valores = [str(v).lower().strip() for v in row.values if str(v) not in ("nan", "")]
    hits = sum(
        1 for keywords in _HEADER_KEYWORDS.values()
        if any(any(kw in v for kw in keywords) for v in valores)
    )
    return hits >= 2


def _montar_registro(
    row_vals: list,
    idx: dict[str, int | None],
    origem_aba: str,
    normalizar_data: bool = True,
) -> dict | None:
    """
    Monta um dicionario de registro a partir de indices resolvidos.

    Args:
        row_vals: Valores da linha.
        idx: Indices das colunas.
        origem_aba: Nome da aba de origem.
        normalizar_data: Se True (padrao), converte a data para string ISO
                         via _normalizar_data. Se False, mantem o valor bruto
                         da celula (usado por estrategias 1-3 para posterior
                         processamento em lote por resolver_datas_ambiguas).
    """
    def _get(campo: str) -> str:
        i = idx.get(campo)
        if i is None or i >= len(row_vals):
            return SENTINEL
        v = row_vals[i]
        try:
            return SENTINEL if pd.isna(v) else str(v).strip()
        except (TypeError, ValueError):
            return str(v).strip()

    data_raw   = _get("data")
    data_str   = _normalizar_data(data_raw) if normalizar_data else data_raw
    valor_str  = _get("valor")
    valor_flt  = _normalizar_valor(valor_str)

    if data_str == SENTINEL and valor_flt == 0.0:
        return None

    return {
        "data":           data_str,
        "conta_contabil": _get("conta_contabil"),
        "valor":          valor_flt,
        "tipo":           _get("tipo"),
        "historico":      _get("historico"),
        "filial":         _get("filial"),
    }


# ---------------------------------------------------------------------------
# ESTRATÉGIA 1 — Semicolon-delimited single cell
# ---------------------------------------------------------------------------
def _estrategia_semicolon(df_raw: pd.DataFrame, nome_aba: str) -> pd.DataFrame | None:
    """
    Detecta e expande arquivos onde todas as colunas estão compactadas
    em uma única célula separada por ';'.

    Retorna DataFrame normalizado ou None se não se aplica.
    """
    # Verifica se ≥80% das linhas têm apenas 1 coluna não-nula com ';' dentro
    col0 = df_raw.iloc[:, 0].dropna().astype(str)
    if df_raw.shape[1] > 2:
        return None
    ratio = (col0.str.contains(";")).mean()
    if ratio < 0.5:
        return None

    logger.info("'%s': Estratégia 1 (semicolon-delimited) ativada.", nome_aba)

    # Expande pela primeira linha (cabeçalho) ou detecta posicionalmente
    primeira = col0.iloc[0]
    tem_cabecalho = _e_linha_cabecalho(pd.Series(primeira.split(";")))

    if tem_cabecalho:
        headers = [h.strip() for h in primeira.split(";")]
        linhas  = col0.iloc[1:]
    else:
        # Infere cabeçalho posicional a partir dos dados
        headers = [f"col_{i}" for i in range(len(primeira.split(";")))]
        linhas  = col0

    registros: list[dict] = []
    for linha in linhas:
        partes = [p.strip() for p in linha.split(";")]
        if len(partes) < len(headers):
            partes += [SENTINEL] * (len(headers) - len(partes))

        if tem_cabecalho:
            idx = {campo: _resolver_coluna(headers, campo) for campo in _COL_ALIASES}
        else:
            # Detecta posições por conteúdo
            idx = _inferir_indices_posicionais(partes, headers)

        reg = _montar_registro(partes, idx, nome_aba, normalizar_data=False)
        if reg:
            registros.append(reg)

    return pd.DataFrame(registros) if registros else None


# ---------------------------------------------------------------------------
# ESTRATÉGIA 2 — Named header detection (multi-bloco)
# ---------------------------------------------------------------------------
def _estrategia_named_header(df_raw: pd.DataFrame, nome_aba: str) -> pd.DataFrame | None:
    """
    Varre todas as linhas buscando cabeçalhos. Extrai um bloco por cabeçalho.
    Retorna DataFrame ou None se nenhum cabeçalho for encontrado.
    """
    header_rows = [i for i, row in df_raw.iterrows() if _e_linha_cabecalho(row)]
    if not header_rows:
        return None

    logger.info("'%s': Estratégia 2 (named header) — cabeçalhos em linhas %s.", nome_aba, header_rows)

    todos_blocos: list[pd.DataFrame] = []

    for bloco_idx, header_row in enumerate(header_rows):
        fim = header_rows[bloco_idx + 1] if bloco_idx + 1 < len(header_rows) else len(df_raw)
        dados = df_raw.iloc[header_row + 1: fim]

        if dados.empty:
            continue

        headers = [str(v).strip() for v in df_raw.iloc[header_row].values]
        idx = {campo: _resolver_coluna(headers, campo) for campo in _COL_ALIASES}

        registros: list[dict] = []
        for _, row in dados.iterrows():
            try:
                reg = _montar_registro(
                    list(row.values), idx, nome_aba, normalizar_data=False,
                )
                if reg:
                    registros.append(reg)
            except Exception as exc:
                logger.warning("'%s' bloco %d linha ignorada: %s", nome_aba, bloco_idx, exc)

        if registros:
            todos_blocos.append(pd.DataFrame(registros))

    return pd.concat(todos_blocos, ignore_index=True) if todos_blocos else None


# ---------------------------------------------------------------------------
# ESTRATÉGIA 3 — Headerless positional extraction
# ---------------------------------------------------------------------------
def _inferir_indices_posicionais(
    amostra_vals: list[str],
    headers: list[str],
) -> dict[str, int | None]:
    """
    Infere índices de colunas por análise de conteúdo quando não há cabeçalho.
    Converte a amostra em DataFrame e delega para _detectar_indices_posicionais.
    """
    df_amostra = pd.DataFrame([amostra_vals], columns=headers)
    return _detectar_indices_posicionais(df_amostra)


def _detectar_indices_posicionais(df_sample: pd.DataFrame) -> dict[str, int | None]:
    """
    Analisa as primeiras linhas do DataFrame para inferir qual coluna
    contem data, valor e conta_contabil — sem depender de nomes.
    """
    idx: dict[str, int | None] = {c: None for c in _COL_ALIASES}
    n_cols = df_sample.shape[1]

    col_scores: dict[str, dict[int, float]] = {
        "data": {}, "valor": {}, "conta_contabil": {}
    }

    for col_i in range(n_cols):
        col_vals = df_sample.iloc[:, col_i].dropna().astype(str).head(20)
        if col_vals.empty:
            continue

        # Score para DATA: % de valores que batem com regex de data
        data_hits = col_vals.apply(
            lambda v: bool(_RE_DATA.search(v) or _RE_DATA_ISO.search(v))
        ).mean()
        col_scores["data"][col_i] = data_hits

        # Score para VALOR: % de valores que são numéricos/monetários
        def _is_numeric(v: str) -> bool:
            v2 = v.replace(".", "").replace(",", "").replace("R$", "").strip()
            return v2.replace("-", "").replace("+", "").isnumeric() and len(v2) > 0

        valor_hits = col_vals.apply(_is_numeric).mean()
        # Penaliza colunas que parecem código (muitos dígitos sem separador decimal)
        media_len = col_vals.str.len().mean()
        if media_len > 8 and valor_hits > 0.8:
            # Pode ser código de conta — verifica se tem separador decimal
            tem_decimal = col_vals.apply(lambda v: "," in v or "." in v).mean()
            if tem_decimal < 0.3:
                valor_hits *= 0.2  # penaliza fortemente
        col_scores["valor"][col_i] = valor_hits

        # Score para CONTA_CONTABIL: texto longo e variado
        media_len_cat = col_vals.str.len().mean()
        n_unicos = col_vals.nunique() / max(len(col_vals), 1)
        col_scores["conta_contabil"][col_i] = (media_len_cat / 100) * n_unicos

    # Atribui melhor coluna por campo (sem repetir índice)
    usados: set[int] = set()
    for campo in ("data", "valor", "conta_contabil"):
        scores = col_scores[campo]
        candidatos = sorted(scores.items(), key=lambda x: x[1], reverse=True)
        for col_i, score in candidatos:
            if col_i not in usados and score > 0.1:
                idx[campo] = col_i
                usados.add(col_i)
                logger.info("Posicional: campo '%s' → coluna %d (score=%.2f)", campo, col_i, score)
                break

    return idx


def _estrategia_posicional(df_raw: pd.DataFrame, nome_aba: str) -> pd.DataFrame | None:
    """
    Extrai dados de arquivos completamente sem cabeçalho,
    inferindo colunas por análise de conteúdo.
    """
    logger.info("'%s': Estratégia 3 (posicional) ativada.", nome_aba)

    idx = _detectar_indices_posicionais(df_raw)

    if idx.get("data") is None and idx.get("valor") is None:
        logger.warning("'%s': posicional não conseguiu inferir nenhum campo.", nome_aba)
        return None

    registros: list[dict] = []
    for _, row in df_raw.iterrows():
        try:
            reg = _montar_registro(
                list(row.values), idx, nome_aba, normalizar_data=False,
            )
            if reg:
                registros.append(reg)
        except Exception as exc:
            logger.warning("'%s' linha posicional ignorada: %s", nome_aba, exc)

    return pd.DataFrame(registros) if registros else None


# ---------------------------------------------------------------------------
# Orquestrador — tenta as 3 estratégias em cascata
# ---------------------------------------------------------------------------
def _processar_aba(df_raw: pd.DataFrame, nome_aba: str) -> pd.DataFrame:
    """
    Tenta as 3 estratégias em ordem. Retorna o resultado da primeira
    que produzir dados, ou DataFrame vazio.
    """
    # Estratégia 1: semicolon em célula única
    resultado = _estrategia_semicolon(df_raw, nome_aba)
    if resultado is not None and not resultado.empty:
        logger.info("'%s': %d registros via Estratégia 1.", nome_aba, len(resultado))
        return resultado

    # Estratégia 2: cabeçalho nomeado
    resultado = _estrategia_named_header(df_raw, nome_aba)
    if resultado is not None and not resultado.empty:
        logger.info("'%s': %d registros via Estratégia 2.", nome_aba, len(resultado))
        return resultado

    # Estratégia 3: posicional sem cabeçalho
    resultado = _estrategia_posicional(df_raw, nome_aba)
    if resultado is not None and not resultado.empty:
        logger.info("'%s': %d registros via Estratégia 3.", nome_aba, len(resultado))
        return resultado

    logger.warning("'%s': nenhuma estratégia extraiu dados.", nome_aba)
    return _df_vazio()


# ---------------------------------------------------------------------------
# ESTRATÉGIA 4 — Fallback via mapeamento fuzzy de colunas
# ---------------------------------------------------------------------------
def _processar_aba_mapeada(
    df_raw: pd.DataFrame, nome_aba: str,
) -> pd.DataFrame:
    """Processa uma unica aba usando mapeamento fuzzy de colunas."""
    if df_raw.empty:
        return _df_vazio()

    cabecalhos = [str(v).strip() for v in df_raw.iloc[0].values]
    dados = df_raw.iloc[1:].reset_index(drop=True)

    mapeamento = mapear_colunas(cabecalhos)
    if not mapeamento:
        return _df_vazio()

    mapa_idx: dict[int, str] = {}
    for i, h in enumerate(cabecalhos):
        if h in mapeamento:
            mapa_idx[i] = mapeamento[h]

    tem_data = "data" in mapa_idx.values()
    tem_valor = "valor" in mapa_idx.values()
    tem_debito = "debito" in mapa_idx.values()
    tem_credito = "credito" in mapa_idx.values()
    if not tem_data or not (tem_valor or tem_debito or tem_credito):
        return _df_vazio()

    registros: list[dict] = []
    for _, row in dados.iterrows():
        reg: dict[str, str] = {}
        for idx, campo in mapa_idx.items():
            if idx < len(row):
                val = row[idx]
                try:
                    reg[campo] = str(val).strip() if not pd.isna(val) else ""
                except (TypeError, ValueError):
                    reg[campo] = str(val).strip()
        registros.append(reg)

    if not registros:
        return _df_vazio()

    df = pd.DataFrame(registros)

    if not tem_valor and (tem_debito or tem_credito):
        col_debito = "debito" if tem_debito else None
        col_credito = "credito" if tem_credito else None
        df = derivar_valor_tipo_de_debito_credito(df, col_debito, col_credito)

    colunas_validas = {c: c for c in df.columns if c in COLUNAS_PADRAO}
    df = df.rename(columns=colunas_validas)
    df = df[[c for c in COLUNAS_PADRAO if c in df.columns]]
    df = limpar_dataframe(df)

    if not df.empty:
        logger.info("'%s': %d registros via fallback.", nome_aba, len(df))

    return df


def _tentar_mapeamento_fallback(arquivo: BinaryIO, ext: str, nome: str) -> pd.DataFrame:
    """
    Fallback quando as 3 estrategias nao encontram dados.

    Le o arquivo raw (todas as abas se Excel), mapeia colunas por fuzzy
    match (rapidfuzz) e extrai registros. Se detectar debito/credito,
    deriva valor+tipo.
    """
    try:
        arquivo.seek(0)
        todos: list[pd.DataFrame] = []

        if ext == "csv":
            conteudo = arquivo.read()
            if isinstance(conteudo, bytes):
                for encoding in ("utf-8", "latin-1", "cp1252"):
                    try:
                        texto = conteudo.decode(encoding)
                        break
                    except UnicodeDecodeError:
                        continue
                else:
                    texto = conteudo.decode("latin-1", errors="replace")
            else:
                texto = str(conteudo)

            primeira = texto.split("\n")[0]
            sep = ";" if primeira.count(";") > primeira.count(",") else ","
            df_raw = pd.read_csv(
                io.StringIO(texto), sep=sep, dtype=str, header=None,
            )
            bloco = _processar_aba_mapeada(df_raw, nome)
            if not bloco.empty:
                todos.append(bloco)

        else:
            xl = pd.ExcelFile(arquivo)
            for sheet in xl.sheet_names:
                try:
                    df_raw = xl.parse(sheet, header=None, dtype=str)
                    bloco = _processar_aba_mapeada(df_raw, sheet)
                    if not bloco.empty:
                        todos.append(bloco)
                except Exception as exc:
                    logger.warning("Aba '%s' ignorada no fallback: %s", sheet, exc)

        if not todos:
            return _df_vazio()

        df = pd.concat(todos, ignore_index=True)
        logger.info("'%s': %d registros via fallback de mapeamento.", nome, len(df))
        return df

    except Exception as exc:
        logger.warning("Fallback de mapeamento falhou para '%s': %s", nome, exc)
        return _df_vazio()


# ---------------------------------------------------------------------------
# Normalização de colunas — ContaView
# ---------------------------------------------------------------------------
_NORMALIZE_ALIASES: dict[str, list[str]] = {
    "data":          ["data", "vencimento", "date", "dt", "dia", "competencia"],
    "conta_contabil": ["conta contábil", "conta_contabil", "conta", "categoria",
                      "descrição", "descricao", "description", "nome",
                      "classificacao", "classificação"],
    "valor":         ["valor", "vlr", "vl", "value", "amount", "total", "r$"],
    "tipo":          ["tipo", "tp", "c/d", "type", "nat", "natureza"],
    "historico":     ["histórico", "historico", "hist", "descricao", "descrição",
                      "complemento", "observação", "observacao", "obs"],
    "filial":        ["filial", "fil", "unidade", "cc", "centro de custo", "centro de custos"],
}


def normalizar_colunas(df: pd.DataFrame) -> pd.DataFrame:
    df = df.copy()

    mapeamento: dict[str, str] = {}
    for col in df.columns:
        col_lower = str(col).lower().strip()
        for padrao, variacoes in _NORMALIZE_ALIASES.items():
            if col_lower in variacoes:
                mapeamento[col] = padrao
                break
        else:
            for padrao, variacoes in _NORMALIZE_ALIASES.items():
                if any(v in col_lower for v in variacoes):
                    mapeamento[col] = padrao
                    break

    df = df.rename(columns=mapeamento)
    df = df[[c for c in COLUNAS_PADRAO if c in df.columns]]
    return df


# ---------------------------------------------------------------------------
# Limpeza e tipagem — ContaView
# ---------------------------------------------------------------------------
def _limpar_valor_contabil(raw) -> float | None:
    if isinstance(raw, (int, float)):
        return round(float(raw), 2)
    try:
        if pd.isna(raw):
            return None
    except (TypeError, ValueError):
        pass
    texto = str(raw).strip().replace("R$", "").replace(" ", "")
    if not texto or texto in ("nan",):
        return None
    if "," in texto and "." in texto:
        if texto.rfind(",") > texto.rfind("."):
            texto = texto.replace(".", "").replace(",", ".")
        else:
            texto = texto.replace(",", "")
    elif "," in texto:
        texto = texto.replace(",", ".")
    try:
        return round(float(texto), 2)
    except ValueError:
        return None


def limpar_dataframe(df: pd.DataFrame) -> pd.DataFrame:
    df = df.copy()

    if "data" in df.columns:
        # Se a coluna ainda contem objetos date (de resolver_datas_ambiguas),
        # converte para datetime64 sem dayfirst (que causa bug com pandas 3.x).
        # Se já é datetime64 (vindo de outro caminho), apenas valida.
        if not pd.api.types.is_datetime64_any_dtype(df["data"]):
            df["data"] = pd.to_datetime(df["data"], errors="coerce")

    if "valor" in df.columns:
        df["valor"] = df["valor"].apply(_limpar_valor_contabil)
        df = df.dropna(subset=["valor"])

    if "tipo" in df.columns:
        df["tipo"] = df["tipo"].astype(str).str.strip().str.upper()
        df["tipo"] = df["tipo"].apply(lambda x: x if x in ("C", "D") else None)

    if "historico" in df.columns:
        df["historico"] = df["historico"].fillna("").astype(str)
        df["historico"] = df["historico"].replace("NÃO ENCONTRADO", "")
    else:
        df["historico"] = ""

    if "filial" in df.columns:
        df["filial"] = df["filial"].fillna("SEM FILIAL").astype(str)
        df["filial"] = df["filial"].replace("NÃO ENCONTRADO", "SEM FILIAL")
    else:
        df["filial"] = "SEM FILIAL"

    if "data" in df.columns:
        df["periodo"] = df["data"].dt.strftime("%Y-%m")

    df = df.reset_index(drop=True)
    return df


# ---------------------------------------------------------------------------
# Extração de período a partir do cabeçalho bruto do arquivo
# ---------------------------------------------------------------------------
_RE_PERIODO = re.compile(
    r"Per[ií]odo:\s*de\s+(\d{2}/\d{2}/\d{4})\s*a\s+(\d{2}/\d{2}/\d{4})"
)
_RE_PERIODO2 = re.compile(
    r"Per[ií]odo\s+(\d{2}/\d{2}/\d{4})\s*(?:a|at[eé])\s*(\d{2}/\d{2}/\d{4})"
)


def extrair_periodo_do_texto(conteudo: bytes) -> tuple[date, date] | None:
    """Procura 'Período: de DD/MM/AAAA a DD/MM/AAAA' nos primeiros bytes."""
    try:
        texto = conteudo.decode("latin-1", errors="replace")
    except Exception:
        return None
    for padrao in (_RE_PERIODO, _RE_PERIODO2):
        m = padrao.search(texto)
        if m:
            try:
                return (
                    datetime.strptime(m.group(1), "%d/%m/%Y").date(),
                    datetime.strptime(m.group(2), "%d/%m/%Y").date(),
                )
            except ValueError:
                continue
    return None


def extrair_periodo_do_nome_arquivo(nome: str) -> tuple[int, int] | None:
    """
    Extrai (mes, ano) do nome do arquivo.

    Reconhece padroes:
      - Nome do mes por extenso/abreviado + ano: Maio_2026, mai2026, JUNHO_2025
      - MM/AAAA, MM-AAAA, MM_AAAA: 05_2026, 05-2026
      - AAAA/MM, AAAA-MM: 2026-05

    Retorna (mes, ano) ou None se nao reconhecer nenhum padrao.
    Usada como fallback por resolver_datas_ambiguas() quando o arquivo
    nao contem datas inequivocas.
    """
    # Remove extensao do arquivo para analise
    nome_base = nome.rsplit(".", 1)[0] if "." in nome else nome

    # 1. Nome do mes por extenso/abreviado (ex: Maio_2026, mai2026)
    m = _RE_PERIODO_NOME_MES_ANO.search(nome_base)
    if m:
        nome_mes = m.group(1).lower()
        ano = int(m.group(2))
        mes = MESES_PT.get(nome_mes)
        if mes and 2000 <= ano <= 2100:
            return (mes, ano)

    # 2. MM/AAAA ou MM-AAAA ou MM_AAAA (ex: 05_2026)
    m = _RE_PERIODO_NOME_MM_AAAA.search(nome_base)
    if m:
        mes, ano = int(m.group(1)), int(m.group(2))
        if 1 <= mes <= 12 and 2000 <= ano <= 2100:
            return (mes, ano)

    # 3. AAAA-MM ou AAAA_MM (ex: 2026-05)
    m = _RE_PERIODO_NOME_AAAA_MM.search(nome_base)
    if m:
        ano, mes = int(m.group(1)), int(m.group(2))
        if 1 <= mes <= 12 and 2000 <= ano <= 2100:
            return (mes, ano)

    return None


# ---------------------------------------------------------------------------
# Resolução de datas ambíguas — algoritmo de mês dominante
# ---------------------------------------------------------------------------
_RE_DATA_RAW = re.compile(r'(\d{1,2})[\/\-](\d{1,2})[\/\-](\d{2,4})')


def resolver_datas_ambiguas(
    datas_raw: pd.Series,
    nome_arquivo: str = "",
    periodo_extraido: tuple[date, date] | None = None,
) -> tuple[pd.Series, list[str]]:
    """
    Resolve datas ambíguas em uma planilha/aba usando algoritmo de mês dominante.

    Processa TODAS as datas JUNTAS (nao linha a linha) para determinar o formato
    correto quando dia e mes sao ambos ≤ 12.

    ALGORITMO:
    1. Primeira passada — identifica linhas inequívocas (token1 > 12, formato DD/MM).
       Coleta o mês de cada uma dessas linhas.
    2. Determina o "mês dominante" do arquivo: o mês mais frequente entre as linhas
       inequívocas. Se nao houver nenhuma, usa fallbacks (periodo do texto do arquivo,
       nome do arquivo, padrao brasileiro DD/MM).
    3. Segunda passada — para cada linha ambígua (ambos tokens ≤ 12), escolhe a
       interpretação cujo mês resultante coincide com o mês dominante.
    4. Linhas que permanecem ambíguas sao registradas em avisos.

    Esta heurística funciona para qualquer planilha futura com formato de data
    inconsistente — nao é uma correcao pontual para um arquivo especifico.

    Args:
        datas_raw: Series de strings de data brutas (ex: "05/01/2026")
        nome_arquivo: Nome do arquivo (para fallback de periodo)
        periodo_extraido: Periodo extraido do texto do cabecalho

    Returns:
        (series_de_datas, avisos) onde series_de_datas contem objetos date
        e avisos é uma lista de strings de aviso.
    """
    from collections import Counter

    avisos: list[str] = []

    if datas_raw.empty:
        return datas_raw, avisos

    # Preserva indices originais e limpa valores
    indices_validos = datas_raw.dropna().index
    valores = datas_raw.dropna().astype(str).str.strip()

    tokens: list[dict | None] = []
    inequivocos_meses: list[int] = []

    for val in valores:
        m = _RE_DATA_RAW.match(val)
        if m:
            t1, t2, ano = int(m.group(1)), int(m.group(2)), int(m.group(3))
            if ano < 100:
                ano += 2000
            tokens.append({"t1": t1, "t2": t2, "ano": ano})
            # Inequívoco se token1 > 12: formato DD/MM, mes = token2
            if t1 > 12:
                inequivocos_meses.append(t2)
            else:
                inequivocos_meses.append(None)
        else:
            tokens.append(None)
            inequivocos_meses.append(None)

    # --- Determinar mês dominante ---
    mes_dominante: int | None = None
    meses_validos = [m for m in inequivocos_meses if m is not None and 1 <= m <= 12]

    if meses_validos:
        contagem = Counter(meses_validos)
        mes_dominante = contagem.most_common(1)[0][0]
        logger.info(
            "Datas: %d inequivocas, mes dominante=%02d (freq=%d)",
            len(meses_validos), mes_dominante, contagem[mes_dominante],
        )
    else:
        # Fallback 1: período extraído do texto do arquivo
        if periodo_extraido:
            mes_dominante = periodo_extraido[1].month
            logger.info(
                "Datas: nenhuma inequivoca, mes dominante via periodo=%02d",
                mes_dominante,
            )
        else:
            # Fallback 2: período extraído do nome do arquivo
            periodo_nome = extrair_periodo_do_nome_arquivo(nome_arquivo)
            if periodo_nome:
                mes_dominante = periodo_nome[0]
                logger.info(
                    "Datas: nenhuma inequivoca, mes dominante via nome=%02d",
                    mes_dominante,
                )
            else:
                # Fallback 3: assume formato DD/MM (padrão brasileiro)
                # Determinar mes dominante das linhas AMBIGUAS como DD/MM
                datas_ddmm = []
                for t in tokens:
                    if t and t["t1"] <= 12 and t["t2"] <= 12:
                        try:
                            d = date(t["ano"], t["t2"], t["t1"])
                            datas_ddmm.append(d.month)
                        except (ValueError, OverflowError):
                            continue
                if datas_ddmm:
                    contagem_ddmm = Counter(datas_ddmm)
                    mes_dominante = contagem_ddmm.most_common(1)[0][0]
                    logger.info(
                        "Datas: nenhuma inequivoca, mes dominante via DD/MM=%02d",
                        mes_dominante,
                    )
                else:
                    mes_dominante = None
                    avisos.append(
                        "Nao foi possivel determinar o periodo automaticamente. "
                        "Informe o periodo manualmente na tela de importacao."
                    )

    # --- Segunda passada: resolver cada data ---
    resultados: list[date | None] = []

    for t in tokens:
        if t is None:
            resultados.append(None)
            continue

        t1, t2, ano = t["t1"], t["t2"], t["ano"]

        try:
            if t1 > 12:
                # Inequívoco DD/MM: dia=t1, mes=t2
                resultados.append(date(ano, t2, t1))
            elif t2 > 12:
                # Inequívoco MM/DD: mes=t1, dia=t2 (caso raro)
                resultados.append(date(ano, t1, t2))
            elif mes_dominante is not None:
                # Ambíguo: usar mês dominante para decidir
                if t1 == mes_dominante:
                    # mes=t1 (MM/DD), dia=t2
                    resultados.append(date(ano, t1, t2))
                elif t2 == mes_dominante:
                    # dia=t1, mes=t2 (DD/MM)
                    resultados.append(date(ano, t2, t1))
                else:
                    # Nenhum bate com mês dominante — fallback DD/MM
                    resultados.append(date(ano, t2, t1))
                    avisos.append(
                        f"Data ambigua nao resolvida: {t1:02d}/{t2:02d}/{ano} "
                        f"(mes dominante={mes_dominante:02d}) — usada como DD/MM."
                    )
            else:
                # Sem mês dominante — fallback DD/MM
                resultados.append(date(ano, t2, t1))
        except (ValueError, OverflowError):
            resultados.append(None)

    # Constroi Series com os mesmos indices da entrada
    series_resultado = pd.Series(
        resultados,
        index=indices_validos,
        dtype="object",
    )

    # Reindexa para o tamanho original (preenche NaT onde era NaN na entrada)
    series_resultado = series_resultado.reindex(datas_raw.index)

    return series_resultado, avisos


# ---------------------------------------------------------------------------
# Leitura principal — retorna dict com metadados
# ---------------------------------------------------------------------------
def ler_arquivo(arquivo: BinaryIO) -> dict:
    nome: str = getattr(arquivo, "name", "arquivo")
    ext = nome.rsplit(".", 1)[-1].lower() if "." in nome else ""
    avisos: list[str] = []

    if ext not in ("xlsx", "xls", "csv"):
        return {
            "sucesso": False, "df": None,
            "linhas_lidas": 0, "linhas_descartadas": 0,
            "motivo_falha": f"Formato '.{ext}' nao suportado para importacao. Use xlsx ou csv.",
            "avisos": avisos,
        }

    # --- Leitura do cabecalho bruto para extrair periodo ---
    periodo_extraido = None
    try:
        cabecalho_raw = arquivo.read(4096)
        arquivo.seek(0)
        periodo_extraido = extrair_periodo_do_texto(cabecalho_raw)
    except Exception:
        arquivo.seek(0)

    try:
        # --- Deteccao de XML SpreadsheetLegacy para .xls ---
        if ext == "xls":
            cabecalho = arquivo.read(200)
            if detectar_xml_legado(cabecalho):
                arquivo.seek(0)
                conteudo = arquivo.read()
                abas = ler_xml_legado(conteudo)
                todos: list[pd.DataFrame] = []
                for nome_aba, df_raw in abas.items():
                    bloco = _processar_aba(df_raw, nome_aba)
                    if not bloco.empty:
                        todos.append(bloco)
                df_processado = pd.concat(todos, ignore_index=True) if todos else _df_vazio()
            else:
                arquivo.seek(0)
                df_processado = processar_excel_csv(arquivo)
        else:
            df_processado = processar_excel_csv(arquivo)

        # --- Fallback via mapeamento fuzzy de colunas ---
        if df_processado.empty and ext in ("xlsx", "xls", "csv"):
            df_fallback = _tentar_mapeamento_fallback(arquivo, ext, nome)
            if df_fallback is not None and not df_fallback.empty:
                df_processado = df_fallback

        if df_processado.empty:
            return {
                "sucesso": False, "df": None,
                "linhas_lidas": 0, "linhas_descartadas": 0,
                "motivo_falha": "Nenhum registro contabil identificado no arquivo.",
                "avisos": avisos,
            }

        # --- Resolucao de datas ambiguas em lote ---
        # Aplica o algoritmo de mes dominante para resolver datas onde
        # dia e mes sao ambos ≤ 12. Processa TODAS as datas JUNTAS.
        if "data" in df_processado.columns:
            datas_resolvidas, avisos_data = resolver_datas_ambiguas(
                df_processado["data"],
                nome_arquivo=nome,
                periodo_extraido=periodo_extraido,
            )
            df_processado["data"] = datas_resolvidas
            avisos.extend(avisos_data)

            # Se o periodo nao foi determinado (arquivo 100% ambiguo)
            # retorna flag para que a interface peça periodo manual
            qtd_valida = datas_resolvidas.notna().sum()
            if qtd_valida > 0:
                periodo_amostra = [
                    d.strftime("%Y-%m") for d in datas_resolvidas.dropna().iloc[:5]
                ]
                if periodo_amostra:
                    periodo_dominante = periodo_amostra[0]
                    if all(p == periodo_dominante for p in periodo_amostra):
                        pass
            else:
                # Nenhuma data resolvida — precisa de periodo manual
                return {
                    "sucesso": False,
                    "df": df_processado,
                    "periodo_necessario": True,
                    "linhas_lidas": len(df_processado),
                    "linhas_descartadas": 0,
                    "avisos": avisos,
                }

        linhas_lidas = len(df_processado)
        df_normalizado = normalizar_colunas(df_processado)
        df_limpo = limpar_dataframe(df_normalizado)
        linhas_descartadas = linhas_lidas - len(df_limpo)

        # --- Fallback de data via periodo extraido (se coluna data nao existe) ---
        if periodo_extraido and (
            "data" not in df_limpo.columns or df_limpo["data"].isna().all()
        ):
            data_fim = pd.Timestamp(periodo_extraido[1])
            df_limpo["data"] = data_fim
            df_limpo["periodo"] = data_fim.strftime("%Y-%m")
            avisos.append(
                "Data nao encontrada como coluna — usada data do "
                "periodo extraido do cabecalho do arquivo."
            )

        # --- Aviso de tipo nao identificado ---
        if "tipo" in df_limpo.columns:
            qtd_sem_tipo = int(df_limpo["tipo"].isna().sum())
            if qtd_sem_tipo > 0:
                avisos.append(
                    f"Tipo nao identificado para {qtd_sem_tipo} lancamento(s) — "
                    "complete manualmente em Lancamentos, se necessario."
                )

        return {
            "sucesso": True, "df": df_limpo,
            "linhas_lidas": linhas_lidas,
            "linhas_descartadas": linhas_descartadas,
            "motivo_falha": None,
            "avisos": avisos,
        }

    except Exception as exc:
        logger.error("Falha ao processar '%s': %s", nome, exc)
        return {
            "sucesso": False, "df": None,
            "linhas_lidas": 0, "linhas_descartadas": 0,
            "motivo_falha": str(exc),
            "avisos": avisos,
        }


# ---------------------------------------------------------------------------
# Parser Excel / CSV — entrada principal (mantida para compatibilidade)
# ---------------------------------------------------------------------------
def processar_excel_csv(arquivo: BinaryIO) -> pd.DataFrame:
    """
    Parser universal para .xlsx e .csv.
    Aplica automaticamente a estratégia correta por aba/arquivo.

    Suporta:
      - Excel com dados semicolon-delimited em célula única
      - Excel/CSV com cabeçalho em qualquer linha, múltiplos blocos
      - CSV sem cabeçalho, encoding latin-1 ou utf-8
      - Arquivos de qualquer tamanho (processa em chunks se necessário)

    Args:
        arquivo: Objeto de arquivo binário.

    Returns:
        pd.DataFrame: Dados normalizados com colunas padronizadas.
    """
    nome: str = getattr(arquivo, "name", "arquivo")
    todos: list[pd.DataFrame] = []

    try:
        if nome.lower().endswith(".csv"):
            # Tenta utf-8 primeiro, depois latin-1
            conteudo = arquivo.read()
            for encoding in ("utf-8", "latin-1", "cp1252"):
                try:
                    texto = conteudo.decode(encoding)
                    break
                except UnicodeDecodeError:
                    continue
            else:
                texto = conteudo.decode("latin-1", errors="replace")

            # Detecta separador: ';' ou ','
            primeira_linha = texto.split("\n")[0]
            sep = ";" if primeira_linha.count(";") > primeira_linha.count(",") else ","

            df_raw = pd.read_csv(
                io.StringIO(texto),
                sep=sep,
                header=None,
                dtype=str,
                on_bad_lines="skip",
            )
            bloco = _processar_aba(df_raw, nome)
            if not bloco.empty:
                todos.append(bloco)

        else:
            xl = pd.ExcelFile(arquivo, engine="openpyxl")
            for sheet in xl.sheet_names:
                try:
                    df_raw = xl.parse(sheet, header=None, dtype=str)
                    bloco = _processar_aba(df_raw, sheet)
                    if not bloco.empty:
                        todos.append(bloco)
                except Exception as exc:
                    logger.warning("Aba '%s' ignorada: %s", sheet, exc)

    except Exception as exc:
        logger.error("Falha crítica ao ler '%s': %s", nome, exc)
        return _df_vazio()

    if not todos:
        return _df_vazio()

    df_final = pd.concat(todos, ignore_index=True)
    logger.info("'%s': %d registros totais extraídos.", nome, len(df_final))
    return df_final


# ---------------------------------------------------------------------------
# Parser PDF
# ---------------------------------------------------------------------------
def extrair_dados_pdf(arquivo: BinaryIO) -> pd.DataFrame:
    try:
        import pdfplumber
    except ImportError:
        logger.error("pdfplumber não instalado.")
        return _df_vazio()

    nome: str = getattr(arquivo, "name", "arquivo.pdf")
    registros: list[dict] = []

    try:
        with pdfplumber.open(arquivo) as pdf:
            for pg_num, page in enumerate(pdf.pages, start=1):
                tabelas = page.extract_tables()
                for tabela in tabelas:
                    if not tabela or len(tabela) < 2:
                        continue
                    headers = [str(h).lower().strip() if h else "" for h in tabela[0]]
                    idx = {campo: _resolver_coluna(headers, campo) for campo in _COL_ALIASES}
                    for linha in tabela[1:]:
                        try:
                            reg = _montar_registro(list(linha), idx, f"PDF pg.{pg_num}")
                            if reg:
                                registros.append(reg)
                        except Exception as exc:
                            logger.warning("PDF pg.%d linha ignorada: %s", pg_num, exc)

                if not tabelas:
                    texto = page.extract_text() or ""
                    df_txt = pd.DataFrame({"col": texto.splitlines()})
                    bloco = _processar_aba(df_txt, f"PDF pg.{pg_num}")
                    if not bloco.empty:
                        registros.extend(bloco.to_dict("records"))

    except Exception as exc:
        logger.error("Falha ao processar PDF '%s': %s", nome, exc)
        return _df_vazio()

    return pd.DataFrame(registros) if registros else _df_vazio()


# ---------------------------------------------------------------------------
# Parser PPTX
# ---------------------------------------------------------------------------
def extrair_dados_powerpoint(arquivo: BinaryIO) -> pd.DataFrame:
    try:
        from pptx import Presentation
    except ImportError:
        logger.error("python-pptx não instalado.")
        return _df_vazio()

    nome: str = getattr(arquivo, "name", "arquivo.pptx")
    registros: list[dict] = []

    try:
        prs = Presentation(arquivo)
        for slide_num, slide in enumerate(prs.slides, start=1):
            for shape in slide.shapes:
                if shape.has_table:
                    tbl = shape.table
                    headers = [tbl.cell(0, c).text.lower().strip() for c in range(len(tbl.columns))]
                    idx = {campo: _resolver_coluna(headers, campo) for campo in _COL_ALIASES}
                    for row_idx in range(1, len(tbl.rows)):
                        try:
                            vals = [tbl.cell(row_idx, c).text.strip() for c in range(len(tbl.columns))]
                            reg = _montar_registro(vals, idx, f"Slide {slide_num}")
                            if reg:
                                registros.append(reg)
                        except Exception as exc:
                            logger.warning("Slide %d tabela ignorada: %s", slide_num, exc)
                elif shape.has_text_frame:
                    texto = "\n".join(p.text for p in shape.text_frame.paragraphs)
                    df_txt = pd.DataFrame({"col": texto.splitlines()})
                    bloco = _processar_aba(df_txt, f"Slide {slide_num}")
                    if not bloco.empty:
                        registros.extend(bloco.to_dict("records"))
    except Exception as exc:
        logger.error("Falha ao processar PPTX '%s': %s", nome, exc)
        return _df_vazio()

    return pd.DataFrame(registros) if registros else _df_vazio()


# ---------------------------------------------------------------------------
# Parser OCR para Comprovantes (Imagem)
# ---------------------------------------------------------------------------
def extrair_dados_comprovante_imagem(imagem_bytes: bytes) -> dict | None:
    """
    Extrai informações estruturadas (Valor total, Data, Hora) de uma imagem
    de comprovante fiscal usando OCR.
    """
    try:
        import easyocr
    except ImportError:
        logger.error("easyocr não instalado. Impossível ler imagens.")
        return None

    try:
        reader = easyocr.Reader(['pt', 'en'], gpu=False, verbose=False)
        resultados = reader.readtext(imagem_bytes, detail=0)
        texto_completo = " ".join(resultados)

        # 1. Extrair Data
        data_match = _RE_DATA.search(texto_completo) or _RE_DATA_ISO.search(texto_completo)
        data_str = data_match.group(1) if data_match else ""

        # 2. Extrair Hora
        hora_match = re.search(r"\b(\d{2}:\d{2}(?::\d{2})?)\b", texto_completo)
        hora_str = hora_match.group(1) if hora_match else ""

        # 3. Extrair Valor Total (pegamos o maior valor monetário encontrado)
        valores_encontrados = _RE_VALOR.findall(texto_completo)
        maior_valor = 0.0
        for v in valores_encontrados:
            val_flt = _normalizar_valor(v)
            if val_flt > maior_valor:
                maior_valor = val_flt

        return {
            "valor": maior_valor if maior_valor > 0 else None,
            "data": data_str if data_str else None,
            "hora": hora_str if hora_str else None,
            "texto_extraido": texto_completo
        }

    except Exception as exc:
        logger.error("Falha ao processar OCR da imagem: %s", exc)
        return None


# ---------------------------------------------------------------------------
# Dispatcher público
# ---------------------------------------------------------------------------
def processar_arquivo(arquivo: BinaryIO) -> tuple[pd.DataFrame, str]:
    nome: str = getattr(arquivo, "name", "")
    ext = nome.rsplit(".", 1)[-1].lower() if "." in nome else ""

    if ext in ("xlsx", "xls"):
        return processar_excel_csv(arquivo), "excel"
    elif ext == "csv":
        return processar_excel_csv(arquivo), "csv"
    elif ext == "pdf":
        return extrair_dados_pdf(arquivo), "pdf"
    elif ext == "pptx":
        return extrair_dados_powerpoint(arquivo), "pptx"
    elif ext in ("png", "jpg", "jpeg"):
        conteudo = arquivo.read()
        if hasattr(arquivo, "seek"):
            arquivo.seek(0)
        dict_dados = extrair_dados_comprovante_imagem(conteudo)
        if dict_dados:
            reg = {
                "data": _normalizar_data(dict_dados.get("data") or ""),
                "categoria": "Despesa de Comprovante (OCR)",
                "valor": dict_dados.get("valor") or 0.0,
                "tipo": "Saída",
                "status": "Pago",
                "prioridade": "Normal",
                "origem_aba": "Imagem OCR",
            }
            return pd.DataFrame([reg]), "imagem"
        return _df_vazio(), "imagem"
    else:
        raise ValueError(f"Extensão '.{ext}' não suportada. Use: xlsx, csv, pdf, pptx, png, jpg ou jpeg.")
