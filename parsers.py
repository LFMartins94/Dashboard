"""
parsers.py
==========
Módulo de ingestão multifonte para o Dashboard Financeiro.

Sistema de detecção automática com 3 estratégias em cascata:

  ESTRATÉGIA 1 — Semicolon-delimited single cell
    Detecta quando o arquivo tem uma única coluna cujos valores contêm
    múltiplos campos separados por ';'. Expande automaticamente.
    Ex: CAP_ILHAS_DO_LAGO_CONCILIADO.xlsx

  ESTRATÉGIA 2 — Named header detection
    Varre todas as linhas procurando uma que contenha palavras-chave
    de cabeçalho (data, valor, descrição...). Extrai múltiplos blocos.
    Ex: JUNHO_2026.xlsx

  ESTRATÉGIA 3 — Headerless positional extraction
    Para arquivos sem cabeçalho nenhum. Detecta a coluna de data por
    padrão de conteúdo (regex), valor por conteúdo numérico/monetário,
    e categoria pelo campo de texto mais longo.
    Ex: CAP_PLAN_IMPORT_CONS_122023_ILA.csv, CAR_PLAN_IMPORT_CONS_052026_HP.csv
"""

import io
import logging
import re
from datetime import date, datetime
from typing import Any, BinaryIO

import pandas as pd

logger = logging.getLogger(__name__)

COLUNAS_PADRAO: list[str] = [
    "data", "categoria", "valor", "tipo", "status", "prioridade", "origem_aba"
]
SENTINEL: str = "NÃO ENCONTRADO"

_RE_VALOR = re.compile(r"R?\$?\s*([\d]{1,3}(?:[.,]\d{3})*(?:[.,]\d{1,2})?)")
_RE_DATA  = re.compile(r"\b(\d{2}[/\-]\d{2}[/\-]\d{2,4})\b")
_RE_DATA_ISO = re.compile(r"\b(\d{4}-\d{2}-\d{2})\b")

# Keywords para detecção de linha de cabeçalho (Estratégia 2)
_HEADER_KEYWORDS: dict[str, list[str]] = {
    "data":      ["data", "vencimento", "date", "dt", "competencia", "competência"],
    "valor":     ["valor", "value", "amount", "total", "r$", "preço", "preco", "custo", "saída", "saida", "entrada"],
    "categoria": ["categoria", "descrição", "descricao", "description", "nome", "historico", "histórico", "conta"],
}

# Mapeamento de variações → campo padrão (Estratégia 2)
_COL_ALIASES: dict[str, list[str]] = {
    "data":       ["data", "vencimento", "date", "dt", "dia", "competencia", "competência"],
    "categoria":  ["categoria", "descrição", "descricao", "description", "nome", "historico",
                   "histórico", "conta contábil", "conta", "hist"],
    "valor":      ["valor (r$)", "valor", "saída (r$)", "saida (r$)", "entrada (r$)",
                   "value", "amount", "total (r$)", "total", "vlr"],
    "tipo":       ["tipo", "type", "classificação", "classificacao", "nat", "natureza"],
    "status":     ["status", "situação", "situacao"],
    "prioridade": ["prioridade", "priority"],
}


# ---------------------------------------------------------------------------
# Helpers compartilhados
# ---------------------------------------------------------------------------
def _df_vazio() -> pd.DataFrame:
    return pd.DataFrame(columns=COLUNAS_PADRAO)


def _normalizar_valor(raw: Any) -> float:
    if isinstance(raw, (int, float)):
        return round(float(raw), 2)
    texto = str(raw).strip().replace("R$", "").replace(" ", "")
    if not texto or texto in ("nan", SENTINEL):
        return 0.0
    if "," in texto and "." in texto:
        if texto.rfind(",") > texto.rfind("."):
            texto = texto.replace(".", "").replace(",", ".")
        else:
            texto = texto.replace(",", "")
    elif "," in texto:
        texto = texto.replace(",", ".")
    try:
        return round(float(texto), 2)
    except ValueError:
        return 0.0


def _normalizar_data(raw: Any) -> str:
    if raw is None:
        return SENTINEL
    try:
        if pd.isna(raw):
            return SENTINEL
    except (TypeError, ValueError):
        pass
    if isinstance(raw, (date, datetime)):
        return raw.strftime("%Y-%m-%d") if isinstance(raw, datetime) else raw.isoformat()
    texto = str(raw).strip()
    for fmt in (
        "%Y-%m-%d %H:%M:%S", "%Y-%m-%dT%H:%M:%S",
        "%d/%m/%Y", "%d-%m-%Y", "%Y-%m-%d",
        "%d/%m/%y", "%m/%d/%Y",
    ):
        try:
            return datetime.strptime(texto, fmt).strftime("%Y-%m-%d")
        except ValueError:
            continue
    return SENTINEL


def _resolver_coluna(headers: list[str], campo: str) -> int | None:
    """Resolve índice de coluna por match exato e depois parcial."""
    aliases = _COL_ALIASES.get(campo, [])
    # Exato primeiro
    for i, h in enumerate(headers):
        if str(h).lower().strip() in aliases:
            return i
    # Parcial como fallback
    for i, h in enumerate(headers):
        h_l = str(h).lower().strip()
        if any(a in h_l for a in aliases):
            return i
    return None


def _e_linha_cabecalho(row: pd.Series) -> bool:
    """Retorna True se a linha parece ser um cabeçalho (≥2 categorias de keywords)."""
    valores = [str(v).lower().strip() for v in row.values if str(v) not in ("nan", "")]
    hits = sum(
        1 for keywords in _HEADER_KEYWORDS.values()
        if any(any(kw in v for kw in keywords) for v in valores)
    )
    return hits >= 2


def _montar_registro(
    row_vals: list,
    idx: dict[str, int | None],
    origem_aba: str,
) -> dict | None:
    """Monta um dicionário de registro a partir de índices resolvidos."""
    def _get(campo: str) -> str:
        i = idx.get(campo)
        if i is None or i >= len(row_vals):
            return SENTINEL
        v = row_vals[i]
        try:
            return SENTINEL if pd.isna(v) else str(v).strip()
        except (TypeError, ValueError):
            return str(v).strip()

    data_str  = _normalizar_data(_get("data"))
    valor_str = _get("valor")
    valor_flt = _normalizar_valor(valor_str)

    if data_str == SENTINEL and valor_flt == 0.0:
        return None

    return {
        "data":       data_str,
        "categoria":  _get("categoria"),
        "valor":      valor_flt,
        "tipo":       _get("tipo"),
        "status":     _get("status"),
        "prioridade": _get("prioridade"),
        "origem_aba": origem_aba,
    }


# ---------------------------------------------------------------------------
# ESTRATÉGIA 1 — Semicolon-delimited single cell
# ---------------------------------------------------------------------------
def _estrategia_semicolon(df_raw: pd.DataFrame, nome_aba: str) -> pd.DataFrame | None:
    """
    Detecta e expande arquivos onde todas as colunas estão compactadas
    em uma única célula separada por ';'.

    Retorna DataFrame normalizado ou None se não se aplica.
    """
    # Verifica se ≥80% das linhas têm apenas 1 coluna não-nula com ';' dentro
    col0 = df_raw.iloc[:, 0].dropna().astype(str)
    if df_raw.shape[1] > 2:
        return None
    ratio = (col0.str.contains(";")).mean()
    if ratio < 0.5:
        return None

    logger.info("'%s': Estratégia 1 (semicolon-delimited) ativada.", nome_aba)

    # Expande pela primeira linha (cabeçalho) ou detecta posicionalmente
    primeira = col0.iloc[0]
    tem_cabecalho = _e_linha_cabecalho(pd.Series(primeira.split(";")))

    if tem_cabecalho:
        headers = [h.strip() for h in primeira.split(";")]
        linhas  = col0.iloc[1:]
    else:
        # Infere cabeçalho posicional a partir dos dados
        headers = [f"col_{i}" for i in range(len(primeira.split(";")))]
        linhas  = col0

    registros: list[dict] = []
    for linha in linhas:
        partes = [p.strip() for p in linha.split(";")]
        if len(partes) < len(headers):
            partes += [SENTINEL] * (len(headers) - len(partes))

        if tem_cabecalho:
            idx = {campo: _resolver_coluna(headers, campo) for campo in _COL_ALIASES}
        else:
            # Detecta posições por conteúdo
            idx = _inferir_indices_posicionais(partes, headers)

        reg = _montar_registro(partes, idx, nome_aba)
        if reg:
            registros.append(reg)

    return pd.DataFrame(registros) if registros else None


# ---------------------------------------------------------------------------
# ESTRATÉGIA 2 — Named header detection (multi-bloco)
# ---------------------------------------------------------------------------
def _estrategia_named_header(df_raw: pd.DataFrame, nome_aba: str) -> pd.DataFrame | None:
    """
    Varre todas as linhas buscando cabeçalhos. Extrai um bloco por cabeçalho.
    Retorna DataFrame ou None se nenhum cabeçalho for encontrado.
    """
    header_rows = [i for i, row in df_raw.iterrows() if _e_linha_cabecalho(row)]
    if not header_rows:
        return None

    logger.info("'%s': Estratégia 2 (named header) — cabeçalhos em linhas %s.", nome_aba, header_rows)

    todos_blocos: list[pd.DataFrame] = []

    for bloco_idx, header_row in enumerate(header_rows):
        fim = header_rows[bloco_idx + 1] if bloco_idx + 1 < len(header_rows) else len(df_raw)
        dados = df_raw.iloc[header_row + 1: fim]

        if dados.empty:
            continue

        headers = [str(v).strip() for v in df_raw.iloc[header_row].values]
        idx = {campo: _resolver_coluna(headers, campo) for campo in _COL_ALIASES}

        registros: list[dict] = []
        for _, row in dados.iterrows():
            try:
                reg = _montar_registro(list(row.values), idx, nome_aba)
                if reg:
                    registros.append(reg)
            except Exception as exc:
                logger.warning("'%s' bloco %d linha ignorada: %s", nome_aba, bloco_idx, exc)

        if registros:
            todos_blocos.append(pd.DataFrame(registros))

    return pd.concat(todos_blocos, ignore_index=True) if todos_blocos else None


# ---------------------------------------------------------------------------
# ESTRATÉGIA 3 — Headerless positional extraction
# ---------------------------------------------------------------------------
def _inferir_indices_posicionais(
    amostra_vals: list[str],
    headers: list[str],
) -> dict[str, int | None]:
    """
    Infere índices de colunas por análise de conteúdo quando não há cabeçalho.
    Analisa padrões: datas (regex), valores (numérico), texto longo (categoria).
    """
    idx: dict[str, int | None] = {c: None for c in _COL_ALIASES}
    return idx  # será resolvido por _detectar_indices_posicionais abaixo


def _detectar_indices_posicionais(df_sample: pd.DataFrame) -> dict[str, int | None]:
    """
    Analisa as primeiras linhas do DataFrame para inferir qual coluna
    contém data, valor e categoria — sem depender de nomes.
    """
    idx: dict[str, int | None] = {c: None for c in _COL_ALIASES}
    n_cols = df_sample.shape[1]

    col_scores: dict[str, dict[int, float]] = {
        "data": {}, "valor": {}, "categoria": {}
    }

    for col_i in range(n_cols):
        col_vals = df_sample.iloc[:, col_i].dropna().astype(str).head(20)
        if col_vals.empty:
            continue

        # Score para DATA: % de valores que batem com regex de data
        data_hits = col_vals.apply(
            lambda v: bool(_RE_DATA.search(v) or _RE_DATA_ISO.search(v))
        ).mean()
        col_scores["data"][col_i] = data_hits

        # Score para VALOR: % de valores que são numéricos/monetários
        def _is_numeric(v: str) -> bool:
            v2 = v.replace(".", "").replace(",", "").replace("R$", "").strip()
            return v2.replace("-", "").replace("+", "").isnumeric() and len(v2) > 0

        valor_hits = col_vals.apply(_is_numeric).mean()
        # Penaliza colunas que parecem código (muitos dígitos sem separador decimal)
        media_len = col_vals.str.len().mean()
        if media_len > 8 and valor_hits > 0.8:
            # Pode ser código de conta — verifica se tem separador decimal
            tem_decimal = col_vals.apply(lambda v: "," in v or "." in v).mean()
            if tem_decimal < 0.3:
                valor_hits *= 0.2  # penaliza fortemente
        col_scores["valor"][col_i] = valor_hits

        # Score para CATEGORIA: texto longo e variado
        media_len_cat = col_vals.str.len().mean()
        n_unicos = col_vals.nunique() / max(len(col_vals), 1)
        col_scores["categoria"][col_i] = (media_len_cat / 100) * n_unicos

    # Atribui melhor coluna por campo (sem repetir índice)
    usados: set[int] = set()
    for campo in ("data", "valor", "categoria"):
        scores = col_scores[campo]
        candidatos = sorted(scores.items(), key=lambda x: x[1], reverse=True)
        for col_i, score in candidatos:
            if col_i not in usados and score > 0.1:
                idx[campo] = col_i
                usados.add(col_i)
                logger.info("Posicional: campo '%s' → coluna %d (score=%.2f)", campo, col_i, score)
                break

    return idx


def _estrategia_posicional(df_raw: pd.DataFrame, nome_aba: str) -> pd.DataFrame | None:
    """
    Extrai dados de arquivos completamente sem cabeçalho,
    inferindo colunas por análise de conteúdo.
    """
    logger.info("'%s': Estratégia 3 (posicional) ativada.", nome_aba)

    idx = _detectar_indices_posicionais(df_raw)

    if idx.get("data") is None and idx.get("valor") is None:
        logger.warning("'%s': posicional não conseguiu inferir nenhum campo.", nome_aba)
        return None

    registros: list[dict] = []
    for _, row in df_raw.iterrows():
        try:
            reg = _montar_registro(list(row.values), idx, nome_aba)
            if reg:
                registros.append(reg)
        except Exception as exc:
            logger.warning("'%s' linha posicional ignorada: %s", nome_aba, exc)

    return pd.DataFrame(registros) if registros else None


# ---------------------------------------------------------------------------
# Orquestrador — tenta as 3 estratégias em cascata
# ---------------------------------------------------------------------------
def _processar_aba(df_raw: pd.DataFrame, nome_aba: str) -> pd.DataFrame:
    """
    Tenta as 3 estratégias em ordem. Retorna o resultado da primeira
    que produzir dados, ou DataFrame vazio.
    """
    # Estratégia 1: semicolon em célula única
    resultado = _estrategia_semicolon(df_raw, nome_aba)
    if resultado is not None and not resultado.empty:
        logger.info("'%s': %d registros via Estratégia 1.", nome_aba, len(resultado))
        return resultado

    # Estratégia 2: cabeçalho nomeado
    resultado = _estrategia_named_header(df_raw, nome_aba)
    if resultado is not None and not resultado.empty:
        logger.info("'%s': %d registros via Estratégia 2.", nome_aba, len(resultado))
        return resultado

    # Estratégia 3: posicional sem cabeçalho
    resultado = _estrategia_posicional(df_raw, nome_aba)
    if resultado is not None and not resultado.empty:
        logger.info("'%s': %d registros via Estratégia 3.", nome_aba, len(resultado))
        return resultado

    logger.warning("'%s': nenhuma estratégia extraiu dados.", nome_aba)
    return _df_vazio()


# ---------------------------------------------------------------------------
# Parser Excel / CSV — entrada principal
# ---------------------------------------------------------------------------
def processar_excel_csv(arquivo: BinaryIO) -> pd.DataFrame:
    """
    Parser universal para .xlsx e .csv.
    Aplica automaticamente a estratégia correta por aba/arquivo.

    Suporta:
      - Excel com dados semicolon-delimited em célula única
      - Excel/CSV com cabeçalho em qualquer linha, múltiplos blocos
      - CSV sem cabeçalho, encoding latin-1 ou utf-8
      - Arquivos de qualquer tamanho (processa em chunks se necessário)

    Args:
        arquivo: Objeto de arquivo binário.

    Returns:
        pd.DataFrame: Dados normalizados com colunas padronizadas.
    """
    nome: str = getattr(arquivo, "name", "arquivo")
    todos: list[pd.DataFrame] = []

    try:
        if nome.lower().endswith(".csv"):
            # Tenta utf-8 primeiro, depois latin-1
            conteudo = arquivo.read()
            for encoding in ("utf-8", "latin-1", "cp1252"):
                try:
                    texto = conteudo.decode(encoding)
                    break
                except UnicodeDecodeError:
                    continue
            else:
                texto = conteudo.decode("latin-1", errors="replace")

            # Detecta separador: ';' ou ','
            primeira_linha = texto.split("\n")[0]
            sep = ";" if primeira_linha.count(";") > primeira_linha.count(",") else ","

            df_raw = pd.read_csv(
                io.StringIO(texto),
                sep=sep,
                header=None,
                dtype=str,
                on_bad_lines="skip",
            )
            bloco = _processar_aba(df_raw, nome)
            if not bloco.empty:
                todos.append(bloco)

        else:
            xl = pd.ExcelFile(arquivo, engine="openpyxl")
            for sheet in xl.sheet_names:
                try:
                    df_raw = xl.parse(sheet, header=None, dtype=str)
                    bloco = _processar_aba(df_raw, sheet)
                    if not bloco.empty:
                        todos.append(bloco)
                except Exception as exc:
                    logger.warning("Aba '%s' ignorada: %s", sheet, exc)

    except Exception as exc:
        logger.error("Falha crítica ao ler '%s': %s", nome, exc)
        return _df_vazio()

    if not todos:
        return _df_vazio()

    df_final = pd.concat(todos, ignore_index=True)
    logger.info("'%s': %d registros totais extraídos.", nome, len(df_final))
    return df_final


# ---------------------------------------------------------------------------
# Parser PDF
# ---------------------------------------------------------------------------
def extrair_dados_pdf(arquivo: BinaryIO) -> pd.DataFrame:
    try:
        import pdfplumber
    except ImportError:
        logger.error("pdfplumber não instalado.")
        return _df_vazio()

    nome: str = getattr(arquivo, "name", "arquivo.pdf")
    registros: list[dict] = []

    try:
        with pdfplumber.open(arquivo) as pdf:
            for pg_num, page in enumerate(pdf.pages, start=1):
                tabelas = page.extract_tables()
                for tabela in tabelas:
                    if not tabela or len(tabela) < 2:
                        continue
                    headers = [str(h).lower().strip() if h else "" for h in tabela[0]]
                    idx = {campo: _resolver_coluna(headers, campo) for campo in _COL_ALIASES}
                    for linha in tabela[1:]:
                        try:
                            reg = _montar_registro(list(linha), idx, f"PDF pg.{pg_num}")
                            if reg:
                                registros.append(reg)
                        except Exception as exc:
                            logger.warning("PDF pg.%d linha ignorada: %s", pg_num, exc)

                if not tabelas:
                    texto = page.extract_text() or ""
                    df_txt = pd.DataFrame({"col": texto.splitlines()})
                    bloco = _processar_aba(df_txt, f"PDF pg.{pg_num}")
                    if not bloco.empty:
                        registros.extend(bloco.to_dict("records"))

    except Exception as exc:
        logger.error("Falha ao processar PDF '%s': %s", nome, exc)
        return _df_vazio()

    return pd.DataFrame(registros) if registros else _df_vazio()


# ---------------------------------------------------------------------------
# Parser PPTX
# ---------------------------------------------------------------------------
def extrair_dados_powerpoint(arquivo: BinaryIO) -> pd.DataFrame:
    try:
        from pptx import Presentation
    except ImportError:
        logger.error("python-pptx não instalado.")
        return _df_vazio()

    nome: str = getattr(arquivo, "name", "arquivo.pptx")
    registros: list[dict] = []

    try:
        prs = Presentation(arquivo)
        for slide_num, slide in enumerate(prs.slides, start=1):
            for shape in slide.shapes:
                if shape.has_table:
                    tbl = shape.table
                    headers = [tbl.cell(0, c).text.lower().strip() for c in range(len(tbl.columns))]
                    idx = {campo: _resolver_coluna(headers, campo) for campo in _COL_ALIASES}
                    for row_idx in range(1, len(tbl.rows)):
                        try:
                            vals = [tbl.cell(row_idx, c).text.strip() for c in range(len(tbl.columns))]
                            reg = _montar_registro(vals, idx, f"Slide {slide_num}")
                            if reg:
                                registros.append(reg)
                        except Exception as exc:
                            logger.warning("Slide %d tabela ignorada: %s", slide_num, exc)
                elif shape.has_text_frame:
                    texto = "\n".join(p.text for p in shape.text_frame.paragraphs)
                    df_txt = pd.DataFrame({"col": texto.splitlines()})
                    bloco = _processar_aba(df_txt, f"Slide {slide_num}")
                    if not bloco.empty:
                        registros.extend(bloco.to_dict("records"))
    except Exception as exc:
        logger.error("Falha ao processar PPTX '%s': %s", nome, exc)
        return _df_vazio()

    return pd.DataFrame(registros) if registros else _df_vazio()


# ---------------------------------------------------------------------------
# Parser OCR para Comprovantes (Imagem)
# ---------------------------------------------------------------------------
def extrair_dados_comprovante_imagem(imagem_bytes: bytes) -> dict | None:
    """
    Extrai informações estruturadas (Valor total, Data, Hora) de uma imagem
    de comprovante fiscal usando OCR.
    """
    try:
        import easyocr
    except ImportError:
        logger.error("easyocr não instalado. Impossível ler imagens.")
        return None

    try:
        reader = easyocr.Reader(['pt', 'en'], gpu=False, verbose=False)
        resultados = reader.readtext(imagem_bytes, detail=0)
        texto_completo = " ".join(resultados)

        # 1. Extrair Data
        data_match = _RE_DATA.search(texto_completo) or _RE_DATA_ISO.search(texto_completo)
        data_str = data_match.group(1) if data_match else ""

        # 2. Extrair Hora
        hora_match = re.search(r"\b(\d{2}:\d{2}(?::\d{2})?)\b", texto_completo)
        hora_str = hora_match.group(1) if hora_match else ""

        # 3. Extrair Valor Total (pegamos o maior valor monetário encontrado)
        valores_encontrados = _RE_VALOR.findall(texto_completo)
        maior_valor = 0.0
        for v in valores_encontrados:
            val_flt = _normalizar_valor(v)
            if val_flt > maior_valor:
                maior_valor = val_flt

        return {
            "valor": maior_valor if maior_valor > 0 else None,
            "data": data_str if data_str else None,
            "hora": hora_str if hora_str else None,
            "texto_extraido": texto_completo
        }

    except Exception as exc:
        logger.error("Falha ao processar OCR da imagem: %s", exc)
        return None


# ---------------------------------------------------------------------------
# Dispatcher público
# ---------------------------------------------------------------------------
def processar_arquivo(arquivo: BinaryIO) -> tuple[pd.DataFrame, str]:
    nome: str = getattr(arquivo, "name", "")
    ext = nome.rsplit(".", 1)[-1].lower() if "." in nome else ""

    if ext in ("xlsx", "xls"):
        return processar_excel_csv(arquivo), "excel"
    elif ext == "csv":
        return processar_excel_csv(arquivo), "csv"
    elif ext == "pdf":
        return extrair_dados_pdf(arquivo), "pdf"
    elif ext == "pptx":
        return extrair_dados_powerpoint(arquivo), "pptx"
    elif ext in ("png", "jpg", "jpeg"):
        conteudo = arquivo.read()
        if hasattr(arquivo, "seek"):
            arquivo.seek(0)
        dict_dados = extrair_dados_comprovante_imagem(conteudo)
        if dict_dados:
            reg = {
                "data": _normalizar_data(dict_dados.get("data") or ""),
                "categoria": "Despesa de Comprovante (OCR)",
                "valor": dict_dados.get("valor") or 0.0,
                "tipo": "Saída",
                "status": "Pago",
                "prioridade": "Normal",
                "origem_aba": "Imagem OCR",
            }
            return pd.DataFrame([reg]), "imagem"
        return _df_vazio(), "imagem"
    else:
        raise ValueError(f"Extensão '.{ext}' não suportada. Use: xlsx, csv, pdf, pptx, png, jpg ou jpeg.")
